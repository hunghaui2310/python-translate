import unittest

import requests
from django.shortcuts import render
from pptx import Presentation
from rest_framework import viewsets, status
from rest_framework.decorators import api_view
from rest_framework.response import Response
from .serializers import TodoSerializer
from .models import Todo, Files
from .serializers import FilesSerializer
from django.conf import settings
from urllib.parse import unquote
import os
import re

# IS_INVALID_GOOGLE_API = False


# Create your views here.
@api_view(['GET'])
def translate(request, file_name, source, target):
    replace_text(unquote(file_name), source, target)
    return Response(unquote(file_name), status=status.HTTP_200_OK)


@api_view(['GET'])
def download(request, file_name):
    file_path = settings.MEDIA_ROOT + '/translated/' + unquote(file_name)
    file_pointer = open(file_path, 'r')
    response = Response(file_pointer,
                        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename=' + unquote(file_name)
    return response


def replace_text(file_name, source, target):
    file_path = settings.MEDIA_ROOT
    prs = Presentation(file_path + '/origin/' + file_name)
    slides = [slide for slide in prs.slides]
    shapes = []

    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        if shape.has_table:
            tbl = shape.table
            row_count = len(tbl.rows)
            col_count = len(tbl.columns)
            for r in range(0, row_count):
                for c in range(0, col_count):
                    cell = tbl.cell(r, c)
                    paragraphs = cell.text_frame.paragraphs
                    for paragraph in paragraphs:
                        for run in paragraph.runs:
                            if is_need_translate(run.text):
                                if auto_translate(run.text, target, source) != "":
                                    run.text = call_google_api_free(run.text, target, source)
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if is_need_translate(run.text):
                        if auto_translate(run.text, target, source) != "":
                            run.text = call_google_api_free(run.text, target, source)
    # create folder "translated" if not exist
    if not os.path.exists(settings.MEDIA_ROOT + "/translated"):
        os.mkdir(settings.MEDIA_ROOT + "/translated")
    prs.save(file_path + '/translated/' + file_name)


def auto_translate(text_translate, target, source='auto'):
    text_translated = call_google_api_free(text_translate, target, source)
    if text_translated == "":
        text_translated = call_google_api_v2(text_translate, target, source)
    return text_translated

def call_google_api_free(text_translate, target, source='auto'):
    query = {'sl': source, 'tl': target, 'dt': 't', 'client': 'gtx', 'q': text_translate}
    response = requests.get("https://translate.googleapis.com/translate_a/single", query)
    if response.ok:
        parse_response = response.json()
        last_index = len(parse_response) - 1
        if parse_response[last_index][0] == target:
            return ""
        return parse_response[0][0][0]
    else:
        return ""


def call_google_api_v2(text_translate, target, source=''):
    if source == 'auto':
        source = ''
    query = {'source': source, 'target': target, 'q': text_translate, 'key': settings.GOOGLE_API_KEY}
    response = requests.get("https://www.googleapis.com/language/translate/v2", query)
    if response.ok:
        parse_response = response.json().data.translations
        if parse_response[0].detectedSourceLanguage == target:
            return ""
        return parse_response[0].translatedText
    else:
        return text_translate


def is_need_translate(origin_text):
    # check is special character
    if any(char.isalpha() for char in origin_text):
        # check is link
        if is_not_link(origin_text):
            return True
    return False


def is_not_link(origin_text):
    regex = re.compile(
        r'^(?:http|ftp)s?://'  # http:// or https://
        r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+(?:[A-Z]{2,6}\.?|[A-Z0-9-]{2,}\.?)|'  # domain...
        r'localhost|'  # localhost...
        r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})'  # ...or ip
        r'(?::\d+)?'  # optional port
        r'(?:/?|[/?]\S+)$', re.IGNORECASE)
    return re.match(regex, origin_text) is None


class TodoView(viewsets.ModelViewSet):
    serializer_class = TodoSerializer
    queryset = Todo.objects.all()


class FilesViewSet(viewsets.ModelViewSet):
    queryset = Files.objects.all()
    serializer_class = FilesSerializer
